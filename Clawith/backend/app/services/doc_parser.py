"""Unified document reader for PDF, Office, and plaintext formats.

Public API:
    doc_read(path, page_range="", max_chars=50000) -> dict
    doc_extract_tables(path, page_range="") -> dict

All paths are absolute file paths (returned by playwright_browser_download
as `file_id`, or provided directly for uploaded files).
"""

from __future__ import annotations

from pathlib import Path

from app.services.playwright_client import DocParseError


_MAX_CHARS_HARD_CAP = 200_000


def _parse_page_range(page_range: str, page_count: int) -> list[int]:
    """Return a list of 0-based page indices. Empty string = all pages."""
    if not page_range.strip():
        return list(range(page_count))
    pages: set[int] = set()
    for chunk in page_range.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            lo, hi = chunk.split("-", 1)
            try:
                lo_i = int(lo) - 1
                hi_i = int(hi) - 1
            except ValueError:
                raise DocParseError(f"invalid page range: {chunk!r}")
            for i in range(max(0, lo_i), min(page_count - 1, hi_i) + 1):
                pages.add(i)
        else:
            try:
                i = int(chunk) - 1
            except ValueError:
                raise DocParseError(f"invalid page number: {chunk!r}")
            if 0 <= i < page_count:
                pages.add(i)
    return sorted(pages)


def _read_pdf(path: Path, page_range: str) -> tuple[str, int]:
    import pypdf
    try:
        reader = pypdf.PdfReader(str(path))
    except Exception as e:
        raise DocParseError(f"PDF parse failed: {e}")
    count = len(reader.pages)
    if count == 0:
        raise DocParseError("PDF has no pages (file may be corrupt)")
    indices = _parse_page_range(page_range, count)
    parts: list[str] = []
    for i in indices:
        try:
            parts.append(reader.pages[i].extract_text() or "")
        except Exception as e:
            raise DocParseError(f"PDF page {i+1} unreadable: {e}")
    return "\n\n".join(parts), count


def _read_docx(path: Path) -> str:
    from docx import Document
    try:
        doc = Document(str(path))
    except Exception as e:
        raise DocParseError(f"DOCX parse failed: {e}")
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:
        raise DocParseError(f"XLSX parse failed: {e}")
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"### {ws.title}")
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(parts)


def _read_pptx(path: Path) -> tuple[str, int]:
    from pptx import Presentation
    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise DocParseError(f"PPTX parse failed: {e}")
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts), len(prs.slides)


def _read_plain(path: Path, encoding: str = "utf-8") -> str:
    try:
        return path.read_text(encoding=encoding)
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


_EXT_FORMAT = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".md": "md",
    ".markdown": "md",
    ".txt": "txt",
    ".csv": "csv",
}


def doc_read(
    file_id_or_path: str,
    page_range: str = "",
    max_chars: int = 50_000,
) -> dict:
    """Read a document and return {text, truncated, format, page_count, max_chars_applied}.

    Supported formats: pdf, docx, xlsx, pptx, md, txt, csv.
    `max_chars` is capped at 200,000 regardless of caller request.
    """
    path = Path(file_id_or_path)
    if not path.exists():
        raise DocParseError(f"file not found: {file_id_or_path}")

    ext = path.suffix.lower()
    fmt = _EXT_FORMAT.get(ext)
    if fmt is None:
        raise DocParseError(f"unsupported extension: {ext}")

    page_count = 1
    if fmt == "pdf":
        text, page_count = _read_pdf(path, page_range)
    elif fmt == "docx":
        text = _read_docx(path)
    elif fmt == "xlsx":
        text = _read_xlsx(path)
    elif fmt == "pptx":
        text, page_count = _read_pptx(path)
    else:  # md, txt, csv
        text = _read_plain(path)

    applied = min(max_chars, _MAX_CHARS_HARD_CAP)
    truncated = len(text) > applied
    if truncated:
        text = text[:applied]

    return {
        "text": text,
        "truncated": truncated,
        "format": fmt,
        "page_count": page_count,
        "max_chars_applied": applied,
    }


def _tables_xlsx(path: Path) -> list[list[list[str]]]:
    from openpyxl import load_workbook
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:
        raise DocParseError(f"XLSX parse failed: {e}")
    tables = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(["" if c is None else str(c) for c in row])
        if rows:
            tables.append(rows)
    return tables


def _tables_pdf(path: Path, page_range: str) -> list[list[list[str]]]:
    import pdfplumber
    try:
        with pdfplumber.open(str(path)) as pdf:
            count = len(pdf.pages)
            indices = _parse_page_range(page_range, count)
            tables: list[list[list[str]]] = []
            for i in indices:
                try:
                    for t in pdf.pages[i].extract_tables() or []:
                        tables.append([["" if c is None else str(c) for c in row] for row in t])
                except Exception:
                    continue  # skip unreadable pages
            return tables
    except Exception as e:
        raise DocParseError(f"PDF table extraction failed: {e}")


def doc_extract_tables(file_id_or_path: str, page_range: str = "") -> dict:
    """Extract structured tables. Supports xlsx and pdf only."""
    path = Path(file_id_or_path)
    if not path.exists():
        raise DocParseError(f"file not found: {file_id_or_path}")

    ext = path.suffix.lower()
    if ext == ".xlsx":
        tables = _tables_xlsx(path)
    elif ext == ".pdf":
        tables = _tables_pdf(path, page_range)
    else:
        raise DocParseError(f"unsupported for table extraction: {ext}")
    return {"tables": tables, "count": len(tables)}
